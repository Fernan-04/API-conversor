"""Fixtures compartidas de los tests.

Los PDFs son archivos reales de `pdfs/`. Para DOCX/PPTX/XLSX no existen originales
reales en el repo (a diferencia de los PDF, ver SPEC §13.2), así que se generan
fixtures pequeñas al vuelo con las mismas librerías que usan los lectores.
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest

REPO = Path(__file__).resolve().parents[1]
FIXTURES = Path(__file__).resolve().parent / "fixtures"
PDFS = REPO / "pdfs"

# PDF de tabla versionado en el repo (guardián de regresión principal). Vive en
# tests/fixtures/ para que la carpeta pdfs/ quede fuera del control de versiones.
VERBOS = FIXTURES / "verbos-objetivos.pdf"
VERBOS_EXPECTED = FIXTURES / "verbos-objetivos.expected.md"

# PDFs de calibración del filtro de tablas (§5). No se versionan (son documentos
# del usuario); si están en pdfs/ local, los tests que los usan corren.
RUBRICA = PDFS / "rubrica1.pdf"
APF1 = PDFS / "APF1_INDICACION.pdf"

requires_rubrica = pytest.mark.skipif(
    not RUBRICA.exists(), reason="falta pdfs/rubrica1.pdf (PDF de calibración)"
)
requires_apf1 = pytest.mark.skipif(
    not APF1.exists(), reason="falta pdfs/APF1_INDICACION.pdf (PDF de calibración)"
)


def _add_hyperlink(paragraph, text: str, url: str) -> None:
    """Inserta un hipervínculo externo en el párrafo (python-docx no tiene API)."""
    from docx.oxml.ns import qn
    from docx.oxml.shared import OxmlElement

    r_id = paragraph.part.relate_to(
        url,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    run = OxmlElement("w:r")
    run.append(OxmlElement("w:rPr"))
    t = OxmlElement("w:t")
    t.text = text
    run.append(t)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


@pytest.fixture
def docx_bytes() -> bytes:
    """Un .docx con título, párrafos (negrita/cursiva inline, enlace), listas y tabla."""
    from docx import Document

    doc = Document()
    doc.add_heading("Título Principal", level=1)
    doc.add_heading("Sección Uno", level=2)
    doc.add_paragraph("Este es un párrafo normal de introducción.")
    p = doc.add_paragraph()
    run = p.add_run("Etiqueta importante")
    run.bold = True

    # Párrafo con formato inline PARCIAL (negrita y cursiva en medio del texto).
    mix = doc.add_paragraph()
    mix.add_run("Texto con ")
    mix.add_run("negrita").bold = True
    mix.add_run(" y ")
    mix.add_run("cursiva").italic = True
    mix.add_run(".")

    # Párrafo con un hipervínculo.
    link_p = doc.add_paragraph()
    link_p.add_run("Ver ")
    _add_hyperlink(link_p, "nuestra web", "https://example.com/x")

    doc.add_paragraph("Primer ítem", style="List Bullet")
    doc.add_paragraph("Segundo ítem", style="List Bullet")

    # Lista numerada.
    doc.add_paragraph("Paso uno", style="List Number")
    doc.add_paragraph("Paso dos", style="List Number")

    table = doc.add_table(rows=2, cols=3)
    table.rows[0].cells[0].text = "A"
    table.rows[0].cells[1].text = "B"
    table.rows[0].cells[2].text = "C"
    table.rows[1].cells[0].text = "1"
    table.rows[1].cells[1].text = "2"
    table.rows[1].cells[2].text = "3"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    """Un .pptx de una diapositiva con título, viñetas, tabla y notas."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = "Diapositiva de Prueba"
    body = slide.placeholders[1].text_frame
    body.text = "Punto uno"
    body.add_paragraph().text = "Punto dos"

    rows, cols = 2, 2
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(4), Inches(5), Inches(1)
    )
    tbl = tbl_shape.table
    tbl.cell(0, 0).text = "Col1"
    tbl.cell(0, 1).text = "Col2"
    tbl.cell(1, 0).text = "x"
    tbl.cell(1, 1).text = "y"

    slide.notes_slide.notes_text_frame.text = "Nota del orador."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def xlsx_bytes() -> bytes:
    """Un .xlsx con dos hojas de datos."""
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Ventas"
    ws1.append(["Producto", "Cantidad", "Precio"])
    ws1.append(["Manzana", 10, 5])
    ws1.append(["Pera", 3, 8])
    ws2 = wb.create_sheet("Resumen")
    ws2.append(["Total", 13])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
