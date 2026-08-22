"""Lector PPTX: implementa el puerto `DocumentReader` (Fase 1).

Usa `python-pptx`. Cada diapositiva es una sección. Mapea:

  - placeholder de título        -> `Heading(1)`
  - texto con >=2 párrafos        -> `ListBlock` (viñetas)
  - texto de un solo párrafo      -> `Paragraph`
  - tablas                        -> `Table`
  - notas del orador              -> `Heading("Notas")` + `Paragraph`s al final
                                     de la sección (bloque aparte, plan §Fase1)
"""

from __future__ import annotations

import io

from doc2md.config import Config
from doc2md.domain.errors import CorruptFileError
from doc2md.domain.models import Document, Element, Heading, ListBlock, Paragraph, Table
from doc2md.text_utils import clean_text


def _text_frame_paragraphs(text_frame, config: Config) -> list[str]:
    out: list[str] = []
    for p in text_frame.paragraphs:
        text = clean_text(p.text, config).strip()
        if text:
            out.append(text)
    return out


def _table_rows(shape, config: Config) -> list[list[str]]:
    table = shape.table
    return [
        [clean_text(cell.text, config) for cell in row.cells]
        for row in table.rows
    ]


class PptxReader:
    """Adaptador de lectura PPTX (puerto `DocumentReader`)."""

    def read(self, data: bytes, filename: str, config: Config) -> Document:
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(data))
        except Exception as exc:  # noqa: BLE001
            raise CorruptFileError(detail=str(exc)) from exc

        sections: list[list[Element]] = []

        for slide in prs.slides:
            elements: list[Element] = []
            title_shape = slide.shapes.title
            # OJO: python-pptx crea un wrapper Python nuevo en cada acceso, así
            # que `shape is title_shape` NO funciona; hay que comparar por id.
            title_id = title_shape.shape_id if title_shape is not None else None

            if title_shape is not None and title_shape.has_text_frame:
                title = clean_text(title_shape.text, config).strip()
                if title:
                    elements.append(Heading(level=1, text=title))

            for shape in slide.shapes:
                if title_id is not None and shape.shape_id == title_id:
                    continue
                if shape.has_table:
                    rows = _table_rows(shape, config)
                    if rows:
                        elements.append(Table(rows=rows))
                elif shape.has_text_frame:
                    paras = _text_frame_paragraphs(shape.text_frame, config)
                    if len(paras) >= 2:
                        elements.append(ListBlock(items=paras))
                    elif paras:
                        elements.append(Paragraph(text=paras[0]))

            # Notas del orador como bloque aparte al final de la sección.
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                note_lines = [
                    clean_text(line, config).strip()
                    for line in (notes_tf.text or "").splitlines()
                    if line.strip()
                ]
                if note_lines:
                    elements.append(Heading(level=min(2, config.heading_max_levels),
                                            text="Notas"))
                    elements.extend(Paragraph(text=l) for l in note_lines)

            sections.append(elements)

        return Document(sections=sections)
